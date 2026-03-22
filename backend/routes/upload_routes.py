from fastapi import APIRouter, UploadFile, File, HTTPException
import uuid
import os
import io

from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import io
from bs4 import BeautifulSoup

from backend.database.chroma_client import collection

router = APIRouter()


def extract_pdf(file_bytes):
    reader = PdfReader(io.BytesIO(file_bytes))
    text = ""

    for i, page in enumerate(reader.pages):
        page_text = page.extract_text()
        # print(f"Page {i} text:", page_text)

        if page_text:
            text += page_text

    return text

def extract_html(file_bytes):

    html = file_bytes.decode("utf-8")

    soup = BeautifulSoup(html, "html.parser")

    return soup.get_text()


def extract_docx(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    text = ""

    for para in doc.paragraphs:
        text += para.text + "\n"

    return text


def extract_txt(file_bytes):
    return file_bytes.decode("utf-8")

def extract_text(file_bytes, filename):

    ext = os.path.splitext(filename)[1].lower()

    if ext == ".pdf":
        return extract_pdf(file_bytes)

    elif ext == ".docx":
        return extract_docx(file_bytes)

    elif ext == ".txt":
        return extract_txt(file_bytes)

    elif ext == ".pptx":
        return extract_pptx(file_bytes)

    elif ext == ".html":
        return extract_html(file_bytes)

    else:
        raise ValueError(f"Unsupported file type: {ext}")
    
def extract_pptx(file_bytes):

    prs = Presentation(io.BytesIO(file_bytes))
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text

def extract_md(file_bytes):

    return file_bytes.decode("utf-8")

def extract_code(file_bytes):

    return file_bytes.decode("utf-8")



def chunk_text(text, chunk_size=1000, overlap=200):
    chunks = []
    start = 0
    text_length = len(text)

    while start < text_length:
        end = min(start + chunk_size, text_length)
        chunks.append(text[start:end])
        start += chunk_size - overlap

    return chunks


@router.post("/upload")
async def upload_file(file: UploadFile = File(...)):

    file_bytes = await file.read()

    filename = file.filename
    ext = os.path.splitext(filename)[1].lower()

    text = extract_text(file_bytes, filename)


    chunks = chunk_text(text)
    document_id = str(uuid.uuid4())

    ids = []
    documents = []
    metadatas = []

    for i, chunk in enumerate(chunks):
        ids.append(str(uuid.uuid4()))
        documents.append(chunk)
        metadatas.append({
            "document_id": document_id,
            "source": filename,
            "chunk": i
        })

    collection.add(
        documents=documents,
        ids=ids,
        metadatas=metadatas
    )

    return {
        "message": "Document stored in ChromaDB",
        "document_id": document_id,
        "filename": filename,
        "chunks_stored": len(chunks)
    }
    
@router.get("/documents")
def list_documents():
    results = collection.get()

    documents = {}

    for meta in results["metadatas"]:
        doc_id = meta["document_id"]
        source = meta["source"]

        if doc_id not in documents:
            documents[doc_id] = {
                "document_id": doc_id,
                "source": source,
                "chunks": 0
            }

        documents[doc_id]["chunks"] += 1

    return list(documents.values())



@router.get("/documents/{document_id}")
def get_document(document_id: str):

    results = collection.get(
        where={"document_id": document_id}
    )

    if not results or not results.get("documents"):
        raise HTTPException(status_code=404, detail="Document not found")

    return {
        "document_id": document_id,
        "total_chunks": len(results["documents"]),
        "documents": results["documents"],
        "metadata": results["metadatas"]
    }