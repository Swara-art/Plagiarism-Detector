import re
import os
import io
import unicodedata
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup

def normalize_text(text: str) -> str:
    text = unicodedata.normalize("NFKC", text)
    text = text.lower()
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def preprocess_for_embedding(text: str) -> str:
    text = normalize_text(text)
    text = re.sub(r'http\S+|www\.\S+', '', text)
    text = re.sub(r'[^\w\s\.\,\!\?\;\:]', ' ', text)
    return re.sub(r'\s+', ' ', text).strip()

def chunk_with_metadata(text: str, chunk_size: int = 500, overlap: int = 100):
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append({"text": text[start:end], "start": start, "end": end})
        start += chunk_size - overlap
    return chunks

# --- File Text Extraction Utilities ---

def extract_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text

def extract_html(file_bytes: bytes) -> str:
    html = file_bytes.decode("utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")
    for script in soup(["script", "style", "nav", "footer", "header"]):
        script.decompose()
    return soup.get_text()

def extract_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_txt(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="ignore")

def extract_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_code(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="ignore")

def extract_text(file_bytes: bytes, filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        return extract_pdf(file_bytes)
    elif ext == ".docx":
        return extract_docx(file_bytes)
    elif ext == ".txt":
        return extract_txt(file_bytes)
    elif ext == ".pptx":
        return extract_pptx(file_bytes)
    elif ext == ".html":
        return extract_html(file_bytes)
    elif ext in (".py", ".js", ".java", ".c", ".cpp", ".ts", ".go", ".rs", ".rust"):
        return extract_code(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {ext}")